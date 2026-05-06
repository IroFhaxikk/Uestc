from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 颜色主题 ──────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # 深蓝背景
C_ACCENT    = RGBColor(0x00, 0xC8, 0xFF)   # 青蓝强调
C_ACCENT2   = RGBColor(0x7C, 0x3A, 0xFF)   # 紫色强调
C_ACCENT3   = RGBColor(0x00, 0xE5, 0x96)   # 绿色强调
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xC8, 0xD8, 0xE8)
C_GOLD      = RGBColor(0xFF, 0xD7, 0x00)
C_CARD1     = RGBColor(0x10, 0x28, 0x45)   # 卡片背景1
C_CARD2     = RGBColor(0x1A, 0x10, 0x40)   # 卡片背景2
C_CARD3     = RGBColor(0x0A, 0x2A, 0x20)   # 卡片背景3

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 工具函数 ──────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_color, alpha=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return txb

def add_para(tf, text, font_size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return p

# ══════════════════════════════════════════════════════════
# 幻灯片 1  封面
# ══════════════════════════════════════════════════════════
slide_layout = prs.slide_layouts[6]  # blank

def make_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ─── Slide 1: 封面 ────────────────────────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
# 装饰条
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT)
add_rect(sl, 0, 7.42, 13.33, 0.08, C_ACCENT2)
# 渐变色块
add_rect(sl, 0, 1.2, 13.33, 0.005, C_ACCENT, line_color=C_ACCENT, line_width=0.5)

add_text(sl, "AI for Code-based Cryptography", 0.6, 0.5, 12, 0.7,
         font_size=15, color=C_ACCENT, bold=False, align=PP_ALIGN.LEFT)

add_text(sl, "融合方案的突出优势", 0.6, 1.3, 12, 1.2,
         font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

add_text(sl, "第四章  AI × 代数的协同优势分析", 0.6, 2.6, 12, 0.6,
         font_size=22, color=C_ACCENT, bold=False, align=PP_ALIGN.LEFT)

add_text(sl,
    "基于论文《AI for Code-based Cryptography》\n"
    "Mohamed Malhou, Ludovic Perret, Kristin Lauter   |   SAC 2025",
    0.6, 3.4, 12, 0.9, font_size=14, color=C_LIGHT, align=PP_ALIGN.LEFT)

# 三个小标签
for i, (label, col) in enumerate([
    ("4.1  相比纯代数方法", C_ACCENT),
    ("4.2  相比纯AI方法",   C_ACCENT2),
    ("4.3  融合的独特价值", C_ACCENT3),
]):
    x = 0.6 + i * 4.2
    add_rect(sl, x, 4.6, 3.8, 0.55, RGBColor(0x10,0x28,0x45),
             line_color=col, line_width=1.5)
    add_text(sl, label, x+0.15, 4.65, 3.5, 0.45,
             font_size=15, color=col, bold=True, align=PP_ALIGN.LEFT)

add_text(sl, "推理速度 < 1 ms  ·  参数覆盖更广  ·  可解释性增强  ·  双向反馈闭环",
         0.6, 5.6, 12, 0.5, font_size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# ─── Slide 2: 目录 ───────────────────────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT)
add_text(sl, "本章结构", 0.5, 0.15, 12, 0.6,
         font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_rect(sl, 0.5, 0.82, 1.5, 0.04, C_ACCENT)

items = [
    ("4.1", "相比纯代数方法的优势", C_ACCENT,
     "无需显式推导不变量 · 更广参数覆盖 · 毫秒级推理 · 可探索性"),
    ("4.2", "相比纯AI方法的优势",   C_ACCENT2,
     "提升泛化能力 · 增强可解释性 · 弥补数据效率 · 理论保障"),
    ("4.3", "融合的独特价值",       C_ACCENT3,
     "双向反馈闭环：代数指导AI，AI发现驱动代数理论"),
]
for i, (num, title, col, sub) in enumerate(items):
    y = 1.2 + i * 1.8
    add_rect(sl, 0.5, y, 12.3, 1.5, C_CARD1, line_color=col, line_width=1.5)
    add_text(sl, num, 0.65, y+0.1, 0.8, 0.8,
             font_size=30, bold=True, color=col, align=PP_ALIGN.LEFT)
    add_rect(sl, 1.4, y+0.15, 0.04, 1.2, col)
    add_text(sl, title, 1.6, y+0.1, 8, 0.55,
             font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(sl, sub,  1.6, y+0.75, 10.5, 0.55,
             font_size=13, color=C_LIGHT, align=PP_ALIGN.LEFT)

# ─── Slide 3: 4.1 相比纯代数方法 ──────────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT)
add_rect(sl, 0, 0, 0.3, 7.5, C_CARD1)
add_rect(sl, 0.3, 0, 0.06, 7.5, C_ACCENT)

add_text(sl, "4.1  相比纯代数方法的优势", 0.6, 0.12, 12, 0.6,
         font_size=24, bold=True, color=C_WHITE)
add_rect(sl, 0.6, 0.75, 5, 0.04, C_ACCENT)

cards = [
    ("① 无需显式推导不变量",
     "代数方法（FGOPT / CMT / syzygy）需通过\nEagon-Northcott 复形、分次 Betti 数等深刻\n代数几何理论手工设计区分不变量\n\nDeepDistinguisher 自动发现隐含结构，\n降低对人类数学洞察力的依赖",
     C_ACCENT, C_CARD1),
    ("② 更广的参数覆盖",
     "FGOPT 仅在高码率有效（R≈1）\nCMT 区分范围 R ∈ [2/3, 1]\nsyzygy 理论上渐近独立于码率，\n但固定参数仍受限\n\nDeepDistinguisher 在中低码率仍有效：\n首次对 MDPC / QC-MDPC 码给出区分结果\n（Table 4，w=10 准确率 97.14%）",
     C_ACCENT, C_CARD1),
    ("③ 推理效率优势",
     "代数区分器（Gröbner 基 + syzygy）\n实际计算需数小时至数天\n\nCMT 复杂度：\nO(((tm/2−k+1)·C(tm/2+d−1,d))^ω)\n\nDeepDistinguisher 推理时间：\n• CPU 单核 < 100 ms\n• GPU 推理 < 1 ms\n（论文 Section 5.1 Inference Complexity）",
     C_ACCENT, C_CARD1),
    ("④ 可探索性",
     "AI 在未知参数区域的成功为代数研究\n提供线索，提示可能存在尚未发现的\n代数不变量或结构特征\n\n示例：puncturing 技术将 n=128 模型\n迁移到 n=1024（m=10, t=2），\n仍获得 70% 准确率",
     C_ACCENT, C_CARD1),
]

positions = [(0.6,1.0,5.9,5.9),(6.8,1.0,5.9,5.9),
             (0.6,4.1,5.9,3.2),(6.8,4.1,5.9,3.2)]
# 2×2 布局
for idx, ((num, body, col, bg), (x,y,w,h)) in enumerate(zip(cards, [
    (0.55, 1.0, 5.9, 2.95),
    (6.85, 1.0, 5.9, 2.95),
    (0.55, 4.15, 5.9, 2.95),
    (6.85, 4.15, 5.9, 2.95),
])):
    add_rect(sl, x, y, w, h, bg, line_color=col, line_width=1.2)
    add_rect(sl, x, y, w, 0.45, col)
    add_text(sl, num, x+0.15, y+0.05, w-0.3, 0.38,
             font_size=14, bold=True, color=C_BG, align=PP_ALIGN.LEFT)
    add_text(sl, body, x+0.15, y+0.55, w-0.3, h-0.65,
             font_size=12, color=C_LIGHT, align=PP_ALIGN.LEFT)

# ─── Slide 4: 实验数据佐证（4.1 数据）────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT)

add_text(sl, "4.1  实验数据佐证：参数覆盖 & 推理效率", 0.5, 0.12, 12.5, 0.6,
         font_size=22, bold=True, color=C_WHITE)
add_rect(sl, 0.5, 0.75, 6, 0.04, C_ACCENT)

# 左：参数覆盖对比表
add_rect(sl, 0.4, 0.92, 6.0, 5.8, C_CARD1, line_color=C_ACCENT, line_width=1)
add_text(sl, "代数区分器参数限制 vs. DeepDistinguisher",
         0.55, 0.97, 5.7, 0.45, font_size=13, bold=True, color=C_ACCENT)

headers = ["方法", "有效码率范围", "MDPC支持"]
col_w   = [1.6, 2.5, 1.3]
col_x   = [0.5, 2.15, 4.65]
row_y   = [1.5, 1.95, 2.4, 2.85, 3.3]
rows = [
    ("FGOPT",          "R ≈ 1（极高码率）",     "✗"),
    ("CMT",            "R ∈ [2/3, 1]",           "✗"),
    ("syzygy",         "渐近无限制，但固定参数受限", "✗"),
    ("DeepDistinguisher","全范围（含低码率）",    "✓"),
]
add_rect(sl, 0.5, 1.45, 5.9, 0.42, RGBColor(0x08,0x1E,0x3C))
for j, h in enumerate(headers):
    add_text(sl, h, col_x[j]+0.05, 1.48, col_w[j], 0.36,
             font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
for i, row in enumerate(rows):
    bg = RGBColor(0x12,0x28,0x42) if i%2==0 else RGBColor(0x0D,0x1E,0x36)
    add_rect(sl, 0.5, row_y[i], 5.9, 0.42, bg)
    col_color = C_ACCENT3 if row[0]=="DeepDistinguisher" else C_LIGHT
    for j, cell in enumerate(row):
        add_text(sl, cell, col_x[j]+0.05, row_y[i]+0.05, col_w[j], 0.36,
                 font_size=11, color=col_color, align=PP_ALIGN.LEFT)

# 高亮 DeepDistinguisher 行
add_rect(sl, 0.5, 3.3, 5.9, 0.42, RGBColor(0x00,0x28,0x18),
         line_color=C_ACCENT3, line_width=1)
for j, cell in enumerate(rows[3]):
    add_text(sl, cell, col_x[j]+0.05, 3.35, col_w[j], 0.36,
             font_size=11, bold=True, color=C_ACCENT3, align=PP_ALIGN.LEFT)

# MDPC 结果
add_rect(sl, 0.5, 3.9, 5.9, 1.7, RGBColor(0x0A,0x22,0x18),
         line_color=C_ACCENT3, line_width=1)
add_text(sl, "首次 MDPC 区分结果（Table 4）",
         0.65, 3.95, 5.6, 0.4, font_size=12, bold=True, color=C_ACCENT3)
mdpc_data = [
    ("MDPC  w=10",  "97.14%", C_ACCENT3),
    ("MDPC  w=14",  "54.90%", C_LIGHT),
    ("QC-MDPC w=6", "98.02%", C_ACCENT3),
]
for i, (label, acc, col) in enumerate(mdpc_data):
    y = 4.45 + i*0.35
    add_text(sl, f"• {label}：准确率 {acc}", 0.7, y, 5.5, 0.32,
             font_size=12, color=col)

# 右：推理时间对比
add_rect(sl, 6.8, 0.92, 6.0, 3.2, C_CARD1, line_color=C_ACCENT, line_width=1)
add_text(sl, "推理时间对比（Section 5.1）",
         6.95, 0.97, 5.7, 0.45, font_size=13, bold=True, color=C_ACCENT)

time_data = [
    ("Gröbner 基计算",   "数小时～数天",  0.9,  C_LIGHT),
    ("syzygy 区分器",    "亚指数级",      0.6,  C_LIGHT),
    ("CMT 区分器",       "超指数（高t）", 0.4,  C_LIGHT),
    ("DeepDistinguisher\n(CPU 单核)", "< 100 ms", 5.0,   C_ACCENT3),
    ("DeepDistinguisher\n(GPU)",      "< 1 ms",   8.0,   C_ACCENT3),
]
bar_colors = [RGBColor(0x44,0x44,0x88), RGBColor(0x44,0x44,0x88),
              RGBColor(0x44,0x44,0x88), C_ACCENT3, C_ACCENT3]
bar_max_w = 3.5
for i, (label, val, ratio, col) in enumerate(time_data):
    y = 1.55 + i*0.5
    add_text(sl, label, 6.9, y, 2.2, 0.45, font_size=10, color=C_LIGHT)
    bar_w = max(0.15, ratio / 8.0 * bar_max_w)
    add_rect(sl, 9.15, y+0.05, bar_w, 0.32, bar_colors[i])
    add_text(sl, val, 9.2+bar_w, y+0.05, 1.5, 0.32,
             font_size=11, bold=(col==C_ACCENT3), color=col)

# 右下：准确率高亮卡片
add_rect(sl, 6.8, 4.3, 6.0, 2.4, C_CARD1, line_color=C_ACCENT2, line_width=1)
add_text(sl, "Goppa Code 区分精度（Table 1, 论文 p.15）",
         6.95, 4.35, 5.7, 0.4, font_size=12, bold=True, color=C_ACCENT2)

acc_data = [
    ("(q=2,m=6) t=2, n=64", "99.12%"),
    ("(q=2,m=6) t=3, n=64", "98.88%"),
    ("(q=3,m=4) t=2, n=64", "98.25%"),
    ("(q=2,m=7) t=3, n=128","99.48%"),
]
for i, (params, acc) in enumerate(acc_data):
    y = 4.85 + i * 0.38
    add_rect(sl, 6.85, y, 5.9, 0.34,
             RGBColor(0x12,0x10,0x3A) if i%2==0 else RGBColor(0x0E,0x0C,0x30))
    add_text(sl, params, 6.95, y+0.03, 4.0, 0.3, font_size=11, color=C_LIGHT)
    add_text(sl, acc, 10.95, y+0.03, 1.6, 0.3, font_size=12,
             bold=True, color=C_ACCENT3, align=PP_ALIGN.RIGHT)

# ─── Slide 5: 4.2 相比纯AI方法 ──────────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT2)
add_rect(sl, 0, 0, 0.3, 7.5, C_CARD2)
add_rect(sl, 0.3, 0, 0.06, 7.5, C_ACCENT2)

add_text(sl, "4.2  相比纯AI方法的优势", 0.6, 0.12, 12, 0.6,
         font_size=24, bold=True, color=C_WHITE)
add_rect(sl, 0.6, 0.75, 5, 0.04, C_ACCENT2)

cards2 = [
    ("① 提升泛化能力",
     "朴素深度学习：逐行展平生成矩阵\n→ 序列长度 kn，注意力代价 O(k²n²)\n\n代数启发设计：\n• 行序列（Row-Sequence）输入 ✓\n• 系统形式预处理（Standard Form）\n• 角嵌入编码有限域元素\n→ 大幅降低计算复杂度，提升泛化",
     C_ACCENT2, C_CARD2),
    ("② 增强可解释性",
     "代数理论为网络学习提供解释框架\n\n观测现象（论文 Section 7）：\n准确率随 mt 增大而退化\n\n代数解释：\n当所需关系的次数超出网络表达能力时\n区分器失效 —— 网络隐式学习行向量间\n多线性关系，关系次数过高则失效",
     C_ACCENT2, C_CARD2),
    ("③ 弥补数据效率不足",
     "纯AI需海量标注数据\n\n代数先验减少模型需从数据中学习的\n信息量，提高数据利用效率\n\n实验证据：对低 t 参数（t≤4）\n仅需 1K～36K 训练步骤即达 99%+\n（Table 2，n=128, m=7）",
     C_ACCENT2, C_CARD2),
    ("④ 理论保障",
     "代数理论为 AI 方法适用范围提供界限\n\n例：当已知某参数范围内不存在低次\n代数不变量时，可预判 AI 在该范围\n可能失效，避免无效计算\n\n实证：t≥9 时（n=128, m=7）\n准确率≤51%，与代数理论预期一致\n（Table 2，论文 p.16）",
     C_ACCENT2, C_CARD2),
]

for idx, ((num, body, col, bg), (x,y,w,h)) in enumerate(zip(cards2, [
    (0.55, 1.0, 5.9, 2.95),
    (6.85, 1.0, 5.9, 2.95),
    (0.55, 4.15, 5.9, 2.95),
    (6.85, 4.15, 5.9, 2.95),
])):
    add_rect(sl, x, y, w, h, bg, line_color=col, line_width=1.2)
    add_rect(sl, x, y, w, 0.45, col)
    add_text(sl, num, x+0.15, y+0.05, w-0.3, 0.38,
             font_size=14, bold=True, color=C_BG, align=PP_ALIGN.LEFT)
    add_text(sl, body, x+0.15, y+0.55, w-0.3, h-0.65,
             font_size=12, color=C_LIGHT, align=PP_ALIGN.LEFT)

# ─── Slide 6: 实验数据佐证（4.2 数据）────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT2)

add_text(sl, "4.2  实验数据佐证：泛化能力 & 理论保障", 0.5, 0.12, 12.5, 0.6,
         font_size=22, bold=True, color=C_WHITE)
add_rect(sl, 0.5, 0.75, 6, 0.04, C_ACCENT2)

# 左：训练步骤 vs 准确率（体现数据效率）
add_rect(sl, 0.4, 0.92, 6.0, 5.9, C_CARD2, line_color=C_ACCENT2, line_width=1)
add_text(sl, "训练步骤 vs. 准确率（泛化能力）\n(q=2, m=7, n=128, Table 2)",
         0.55, 0.97, 5.7, 0.65, font_size=12, bold=True, color=C_ACCENT2)

step_data = [
    ("t=2", "2K steps",   98.14, 5.0),
    ("t=3", "91K steps",  99.48, 5.0),
    ("t=4", "36K steps",  98.88, 5.0),
    ("t=5", "579K steps", 64.52, 1.65),
    ("t=6", "115K steps", 57.00, 0.7),
    ("t=7", "411K steps", 54.42, 0.45),
]
for i, (label, steps, acc, bar_ratio) in enumerate(step_data):
    y = 1.75 + i * 0.68
    col = C_ACCENT2 if acc >= 90 else (C_LIGHT if acc >= 55 else RGBColor(0x88,0x88,0x88))
    add_text(sl, label, 0.55, y, 0.7, 0.5, font_size=12, bold=True, color=col)
    add_text(sl, steps, 1.3, y, 1.3, 0.5, font_size=11, color=C_LIGHT)
    bar_w = bar_ratio / 5.0 * 2.8
    add_rect(sl, 2.65, y+0.1, bar_w, 0.3, col)
    add_text(sl, f"{acc}%", 2.7+bar_w, y+0.05, 1.0, 0.4,
             font_size=12, bold=(acc>=90), color=col)

add_rect(sl, 0.5, 5.85, 5.9, 0.85, RGBColor(0x14,0x0C,0x38),
         line_color=C_ACCENT2, line_width=1)
add_text(sl, "t ≤ 4：千级步骤达 99%+（高数据效率）\n"
             "t ≥ 9：准确率 ≤ 51%（符合代数理论界限）",
         0.65, 5.9, 5.6, 0.75, font_size=12, color=C_LIGHT)

# 右：架构设计对比
add_rect(sl, 6.8, 0.92, 6.0, 3.0, C_CARD2, line_color=C_ACCENT2, line_width=1)
add_text(sl, "代数启发架构设计 vs. 朴素设计",
         6.95, 0.97, 5.7, 0.45, font_size=12, bold=True, color=C_ACCENT2)

arch_data = [
    ("输入表示", "逐行展平（flat）", "行序列（row-seq）"),
    ("序列长度", "kn", "k（行数）"),
    ("注意力代价", "O(k²n²)", "O(k²)"),
    ("编码方式", "原始整数", "角嵌入/多项式系数"),
    ("Table 6 对比", "Patch 嵌入略差", "Row-seq 最优 ✓"),
]
add_rect(sl, 6.85, 1.5, 5.9, 0.38, RGBColor(0x12,0x06,0x38))
for j, h in enumerate(["对比项", "朴素方案", "代数启发方案"]):
    add_text(sl, h, [6.9, 8.25, 10.5][j], 1.54,
             [1.3, 1.6, 1.9][j], 0.32, font_size=11, bold=True,
             color=C_ACCENT2, align=PP_ALIGN.LEFT)
for i, row in enumerate(arch_data):
    y = 1.95 + i * 0.38
    bg = RGBColor(0x14, 0x08, 0x3C) if i%2==0 else RGBColor(0x10,0x06,0x30)
    add_rect(sl, 6.85, y, 5.9, 0.36, bg)
    for j, cell in enumerate(row):
        col = C_ACCENT3 if j==2 else C_LIGHT
        add_text(sl, cell, [6.9, 8.25, 10.5][j], y+0.03,
                 [1.3, 1.6, 1.9][j], 0.3, font_size=10,
                 color=col, align=PP_ALIGN.LEFT)

# 右下：DeepRecover 展示
add_rect(sl, 6.8, 4.1, 6.0, 2.7, C_CARD2, line_color=C_ACCENT3, line_width=1)
add_text(sl, "DeepRecover：验证代数结构可学习性（Section 6）",
         6.95, 4.15, 5.7, 0.45, font_size=12, bold=True, color=C_ACCENT3)
add_text(sl,
    "任务：给定含缺失项的生成矩阵 G̃ ∈ M(Fq ∪ {*})\n"
    "恢复缺失值使其成为合法 Goppa 码生成矩阵\n\n"
    "实验结果（Table 5，q=2, n=64, m=6）：\n"
    "  t=2：组件准确率 80%\n"
    "  t=3：组件准确率 76%\n"
    "  t=4：组件准确率 64%\n\n"
    "证明：Goppa 码结构可被 AI 学习并利用",
    6.95, 4.65, 5.7, 2.05, font_size=11, color=C_LIGHT)

# ─── Slide 7: 4.3 融合的独特价值 ─────────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT3)
add_rect(sl, 0, 0, 0.3, 7.5, C_CARD3)
add_rect(sl, 0.3, 0, 0.06, 7.5, C_ACCENT3)

add_text(sl, "4.3  融合的独特价值：双向反馈闭环", 0.6, 0.12, 12, 0.6,
         font_size=24, bold=True, color=C_WHITE)
add_rect(sl, 0.6, 0.75, 5.5, 0.04, C_ACCENT3)

# 中央闭环图示（用矩形+箭头文字模拟）
add_rect(sl, 3.5, 1.1, 6.3, 1.05, RGBColor(0x0A,0x26,0x1A),
         line_color=C_ACCENT3, line_width=1.5)
add_text(sl, "代数理论", 3.65, 1.15, 2.8, 0.5,
         font_size=16, bold=True, color=C_ACCENT3, align=PP_ALIGN.LEFT)
add_text(sl, "（FGOPT / CMT / syzygy / Betti 数）",
         3.65, 1.6, 5.8, 0.45, font_size=11, color=C_LIGHT)

add_rect(sl, 3.5, 3.5, 6.3, 1.05, RGBColor(0x0A,0x14,0x38),
         line_color=C_ACCENT, line_width=1.5)
add_text(sl, "深度学习 / AI 模型", 3.65, 3.55, 3.5, 0.5,
         font_size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
add_text(sl, "（DeepDistinguisher / DeepRecover / Transformer）",
         3.65, 3.98, 5.8, 0.45, font_size=11, color=C_LIGHT)

# 箭头（向下 →）
add_text(sl, "指导设计 →\n（输入表示、训练目标）",
         3.65, 2.25, 6.1, 0.9, font_size=12,
         color=C_ACCENT3, italic=True)
# 箭头（向上 ←）
add_text(sl, "← 实验发现驱动\n（新参数区域、隐含不变量线索）",
         3.65, 2.8, 6.1, 0.6, font_size=12,
         color=C_ACCENT, italic=True, align=PP_ALIGN.RIGHT)

# 三个应用场景
scenarios = [
    ("发现新不变量",
     "AI 在 QC-MDPC w=6 区分\n成功（98.02%），提示\n可能存在未被发现的代数\n不变量（论文 Section 5.3）",
     C_ACCENT3, 0.5),
    ("跨码族迁移",
     "puncturing 技术将\nm=7,n=128 模型迁移到\nm=10,n=1024，准确率\n仍达 70%（Section 5.2）",
     C_ACCENT,  4.7),
    ("AI辅助数学发现",
     "\"AI-assisted mathematics\"\n范式：非线性相关性→\n代数结构假设→\n严格数学证明",
     C_ACCENT2, 8.9),
]
for label, body, col, x in scenarios:
    add_rect(sl, x, 4.75, 3.9, 2.55, C_CARD1, line_color=col, line_width=1.2)
    add_rect(sl, x, 4.75, 3.9, 0.42, col)
    add_text(sl, label, x+0.15, 4.78, 3.6, 0.36,
             font_size=13, bold=True, color=C_BG)
    add_text(sl, body,  x+0.15, 5.25, 3.6, 1.95,
             font_size=11.5, color=C_LIGHT)

# ─── Slide 8: 总结对比 ───────────────────────────────────
sl = make_slide()
add_rect(sl, 0, 0, 13.33, 7.5, C_BG)
add_rect(sl, 0, 0, 13.33, 0.08, C_GOLD)

add_text(sl, "融合方案综合优势总结", 0.5, 0.12, 12, 0.6,
         font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_rect(sl, 0.5, 0.75, 4, 0.04, C_GOLD)

# 综合对比表
headers3 = ["评估维度", "纯代数方法", "纯AI方法", "融合方案（本文）"]
col_xs = [0.45, 3.05, 5.55, 8.65]
col_ws = [2.55, 2.4, 3.0, 4.4]
table_rows = [
    ("参数覆盖范围",  "高码率受限",  "无理论指导",  "全范围，首次MDPC结果"),
    ("推理速度",      "小时～天级",  "毫秒级",       "< 1 ms (GPU)  ✓"),
    ("可解释性",      "完全可解释",  "黑盒",         "代数理论提供框架"),
    ("数据需求",      "无需数据",    "海量数据",     "代数先验减少需求"),
    ("理论保障",      "严格理论",    "经验驱动",     "代数界限 + 实验验证"),
    ("新发现能力",    "依赖人类洞察","隐含发现",     "双向反馈，发现新线索"),
]

add_rect(sl, 0.4, 0.92, 12.5, 0.48, RGBColor(0x1A,0x14,0x06))
for j, h in enumerate(headers3):
    col = [C_LIGHT, C_ACCENT, C_ACCENT2, C_GOLD][j]
    add_text(sl, h, col_xs[j]+0.08, 0.97, col_ws[j]-0.1, 0.38,
             font_size=13, bold=True, color=col, align=PP_ALIGN.LEFT)

for i, row in enumerate(table_rows):
    y = 1.48 + i * 0.72
    bg = RGBColor(0x10,0x18,0x24) if i%2==0 else RGBColor(0x0D,0x14,0x20)
    add_rect(sl, 0.4, y, 12.5, 0.68, bg)
    for j, cell in enumerate(row):
        col = [C_LIGHT, C_LIGHT, C_LIGHT, C_GOLD][j]
        if j==3: col = C_ACCENT3
        add_text(sl, cell, col_xs[j]+0.08, y+0.08, col_ws[j]-0.1, 0.55,
                 font_size=11.5, color=col, align=PP_ALIGN.LEFT)

# 底部结论
add_rect(sl, 0.4, 5.95, 12.5, 1.35, RGBColor(0x12,0x12,0x06),
         line_color=C_GOLD, line_width=1.5)
add_text(sl, "核心结论",
         0.6, 5.98, 2.0, 0.45, font_size=14, bold=True, color=C_GOLD)
add_text(sl,
    "融合方案形成了「代数理论 ⇌ AI 模型」双向反馈闭环，兼具代数方法的可解释性与 AI 的参数覆盖广度。\n"
    "DeepDistinguisher 在 CPU 上推理 < 100 ms、GPU < 1 ms，且首次对 MDPC/QC-MDPC 码给出区分结果（w=10, 97.14%），\n"
    "开创了「AI 辅助密码分析」的新范式，为后量子密码安全评估提供了全新工具。",
    0.6, 6.42, 12.1, 0.85, font_size=11.5, color=C_LIGHT)

# ─── 保存 ────────────────────────────────────────────────
out = "/Users/krito/Desktop/研二/纠错编码/融合方案优势分析.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
